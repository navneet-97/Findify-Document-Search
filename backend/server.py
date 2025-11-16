from fastapi import FastAPI, APIRouter, UploadFile, File, Form, HTTPException, Query
from fastapi.responses import FileResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
import uuid
from datetime import datetime, timezone
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
from transformers.pipelines import pipeline
import PyPDF2
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
import io
import shutil

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ.get('MONGO_URL', 'mongodb://localhost:27017/findify')
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ.get('DB_NAME', 'findify')]

# Chroma DB setup
CHROMA_DB_PATH = ROOT_DIR / 'chroma_db'
CHROMA_DB_PATH.mkdir(exist_ok=True, parents=True)
try:
    chroma_client = chromadb.PersistentClient(path=str(CHROMA_DB_PATH))
    collection = chroma_client.get_or_create_collection(
        name="documents",
        metadata={"hnsw:space": "cosine"}
    )
except Exception as e:
    logging.error(f"Error initializing ChromaDB: {e}")
    raise

# Load embedding model
try:
    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')  # Small, efficient model
except Exception as e:
    logging.error(f"Error loading SentenceTransformer: {e}")
    raise

# Load zero-shot classifier for auto-categorization
try:
    classifier = pipeline(
        "zero-shot-classification",
        model="valhalla/distilbart-mnli-12-3",  # Smaller, faster model
        device=-1  # CPU
    )
except Exception as e:
    logging.error(f"Error loading zero-shot classifier: {e}")
    raise

# Category labels for auto-tagging
CANDIDATE_LABELS = [
    "marketing", "sales", "product", "design", "engineering",
    "finance", "hr", "legal", "operations", "research",
    "campaign", "strategy", "analytics", "branding", "content"
]

# Create uploads directory
UPLOAD_DIR = ROOT_DIR / 'uploads'
UPLOAD_DIR.mkdir(exist_ok=True)

# Create the main app without a prefix
app = FastAPI()

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# Add startup event to create indexes
@app.on_event("startup")
async def startup_event():
    # Create indexes for better search performance
    await db.documents.create_index([("file_type", 1)])
    await db.documents.create_index([("tags", 1)])
    await db.documents.create_index([("name", "text"), ("text_content", "text")])
    await db.documents.create_index([("created_at", -1)])

# Define Models
class Document(BaseModel):
    model_config = ConfigDict(extra="ignore")
    
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    file_type: str
    size: int
    tags: List[str] = []
    file_path: str
    text_content: str = ""
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class DocumentResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    
    id: str
    name: str
    file_type: str
    size: int
    tags: List[str]
    file_path: str
    snippet: str = ""
    created_at: str
    updated_at: str
    score: Optional[float] = None

class SearchQuery(BaseModel):
    query: str
    file_types: Optional[List[str]] = None
    tags: Optional[List[str]] = None
    date_from: Optional[str] = None
    date_to: Optional[str] = None
    limit: int = 20

class UpdateTagsRequest(BaseModel):
    tags: List[str]

# Helper functions
def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF"""
    try:
        pdf_file = io.BytesIO(file_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logging.error(f"Error extracting PDF text: {e}")
        return ""

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX"""
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()
    except Exception as e:
        logging.error(f"Error extracting DOCX text: {e}")
        return ""

def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from PPTX"""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                # Simple approach to avoid static analysis issues
                try:
                    text += str(getattr(shape, 'text', ''))
                except:
                    # Skip problematic shapes
                    continue
        return text.strip()
    except Exception as e:
        logging.error(f"Error extracting PPTX text: {e}")
        return ""

def extract_text_from_image(file_bytes: bytes) -> str:
    """Extract text from image using OCR"""
    try:
        image = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        logging.error(f"Error extracting image text: {e}")
        return ""

def extract_text(file_bytes: bytes, file_type: str) -> str:
    """Extract text based on file type"""
    if file_type == "application/pdf":
        return extract_text_from_pdf(file_bytes)
    elif file_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
        return extract_text_from_docx(file_bytes)
    elif file_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
        return extract_text_from_pptx(file_bytes)
    elif file_type.startswith("image/"):
        return extract_text_from_image(file_bytes)
    elif file_type == "text/plain":
        return file_bytes.decode('utf-8', errors='ignore')
    return ""

def auto_categorize(text: str) -> List[str]:
    """Auto-categorize document using zero-shot classification"""
    if not text or len(text.strip()) < 10:
        return []
    
    try:
        # Truncate text for classification (take first 500 chars)
        text_sample = text[:500]
        result = classifier(text_sample, CANDIDATE_LABELS, multi_label=True)
        
        # Return tags with score > 0.3
        if isinstance(result, dict) and 'labels' in result and 'scores' in result:
            tags = [label for label, score in zip(result['labels'], result['scores']) if score > 0.3]
            return tags[:5]  # Return top 5 tags
        return []
    except Exception as e:
        logging.error(f"Error in auto-categorization: {e}")
        return []

def get_snippet(text: str, max_length: int = 200) -> str:
    """Get text snippet for preview"""
    if len(text) <= max_length:
        return text
    return text[:max_length] + "..."

# API Routes
@api_router.get("/")
async def root():
    return {"message": "Document Search API"}

@api_router.post("/documents/upload")
async def upload_document(file: UploadFile = File(...)):
    """Upload and process a document"""
    try:
        # Read file
        file_bytes = await file.read()
        file_size = len(file_bytes)
        
        # Check file size (max 20MB)
        if file_size > 20 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="File size exceeds 20MB limit")
        
        # Save file
        doc_id = str(uuid.uuid4())
        filename = file.filename or f"document_{doc_id}"
        file_ext = Path(filename).suffix or ".bin"
        file_path = UPLOAD_DIR / f"{doc_id}{file_ext}"
        
        with open(file_path, "wb") as f:
            f.write(file_bytes)
        
        # Extract text
        text_content = extract_text(file_bytes, file.content_type or "")
        
        # Auto-categorize
        tags = auto_categorize(text_content)
        
        # Generate embeddings
        if text_content:
            embedding = embedding_model.encode(text_content).tolist()
            
            # Store in Chroma
            collection.add(
                ids=[doc_id],
                embeddings=[embedding],
                documents=[text_content],
                metadatas=[{"doc_id": doc_id, "name": filename}]
            )
        
        # Store metadata in MongoDB
        doc = Document(
            id=doc_id,
            name=filename,
            file_type=file.content_type or "unknown",
            size=file_size,
            tags=tags,
            file_path=str(file_path),
            text_content=text_content
        )
        
        doc_dict = doc.model_dump()
        doc_dict['created_at'] = doc_dict['created_at'].isoformat()
        doc_dict['updated_at'] = doc_dict['updated_at'].isoformat()
        
        await db.documents.insert_one(doc_dict)
        
        return {
            "id": doc_id,
            "name": filename,
            "file_type": file.content_type,
            "size": file_size,
            "tags": tags,
            "message": "Document uploaded successfully"
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Error uploading document: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/documents/search")
async def search_documents(query: SearchQuery):
    """Hybrid search: semantic + keyword"""
    try:
        results = []
        
        # Always build the filter regardless of whether there's a query
        mongo_filter = {}
        
        # Apply filters
        if query.file_types:
            # Handle image/ filter specially to match all image types
            image_filter_types = [ft for ft in query.file_types if ft == "image/"]
            other_types = [ft for ft in query.file_types if ft != "image/"]
            
            if image_filter_types:
                # Match all image types
                image_types = ["image/png", "image/jpeg", "image/jpg", "image/gif", "image/webp"]
                if other_types:
                    mongo_filter['$or'] = [
                        {'file_type': {'$in': other_types}},
                        {'file_type': {'$regex': '^image/'}}
                    ]
                else:
                    mongo_filter['file_type'] = {'$regex': '^image/'}
            elif other_types:
                mongo_filter['file_type'] = {'$in': other_types}
        
        if query.tags:
            mongo_filter['tags'] = {'$in': query.tags}
        
        if query.date_from:
            mongo_filter['created_at'] = {'$gte': query.date_from}
        
        if query.date_to:
            if 'created_at' in mongo_filter:
                mongo_filter['created_at']['$lte'] = query.date_to
            else:
                mongo_filter['created_at'] = {'$lte': query.date_to}
        
        # Optimize search by checking if we have any search query
        if query.query and query.query.strip():
            search_term = query.query.strip().lower()
            
            # Build the search filter to prioritize exact matches
            search_filter = {
                '$or': [
                    {'name': {'$regex': f'^.*{search_term}.*$', '$options': 'i'}},  # Partial name match
                    {'tags': {'$in': [search_term]}}  # Exact tag match
                ]
            }
            
            # Combine with existing filters
            final_filter = {'$and': [mongo_filter, search_filter]} if mongo_filter else search_filter
            
            # Query MongoDB with the combined filter
            docs = await db.documents.find(final_filter, {"_id": 0}).sort('created_at', -1).limit(query.limit).to_list(query.limit)
            
            # Process results
            for doc in docs:
                results.append(DocumentResponse(
                    id=doc['id'],
                    name=doc['name'],
                    file_type=doc['file_type'],
                    size=doc['size'],
                    tags=doc['tags'],
                    file_path=doc['file_path'],
                    snippet=get_snippet(doc.get('text_content', '')),
                    created_at=doc['created_at'],
                    updated_at=doc['updated_at']
                ))
            
            # If we don't have enough results and there's no file type filter, try semantic search
            if len(results) < query.limit and not query.file_types:
                try:
                    query_embedding = embedding_model.encode(search_term).tolist()
                    chroma_results = collection.query(
                        query_embeddings=[query_embedding],
                        n_results=min(query.limit - len(results), 50)
                    )
                    
                    # Process semantic search results
                    if chroma_results and chroma_results.get('ids') and len(chroma_results['ids']) > 0 and len(chroma_results['ids'][0]) > 0:
                        ids = chroma_results['ids'][0]
                        distances_data = chroma_results.get('distances')
                        distances = []
                        if distances_data and isinstance(distances_data, list) and len(distances_data) > 0:
                            distances = distances_data[0]
                        
                        # Create a map of document IDs to scores
                        semantic_scores = {}
                        for i, doc_id in enumerate(ids):
                            distance = distances[i] if i < len(distances) else 0
                            # Convert distance to similarity score (0-1)
                            score = 1 / (1 + distance) if distance is not None else 0
                            semantic_scores[doc_id] = score
                        
                        # Fetch documents from MongoDB with scores, excluding already found ones
                        existing_ids = [r.id for r in results]
                        additional_filter = {"id": {"$nin": existing_ids}}
                        if mongo_filter:
                            combined_filter = {"$and": [mongo_filter, additional_filter]}
                        else:
                            combined_filter = additional_filter
                            
                        docs = await db.documents.find({"id": {"$in": ids}, **combined_filter}, {"_id": 0}).to_list(len(ids))
                        
                        # Combine documents with scores
                        semantic_results = []
                        for doc in docs:
                            doc_id = doc['id']
                            score = semantic_scores.get(doc_id, 0)
                            
                            semantic_results.append(DocumentResponse(
                                id=doc_id,
                                name=doc['name'],
                                file_type=doc['file_type'],
                                size=doc['size'],
                                tags=doc['tags'],
                                file_path=doc['file_path'],
                                snippet=get_snippet(doc.get('text_content', '')),
                                created_at=doc['created_at'],
                                updated_at=doc['updated_at'],
                                score=score
                            ))
                        
                        # Sort by score (descending) and add to results
                        semantic_results.sort(key=lambda x: x.score or 0, reverse=True)
                        results.extend(semantic_results[:query.limit - len(results)])
                        
                except Exception as e:
                    logging.error(f"Error in semantic search: {e}")
                    # Fallback to text search in MongoDB if semantic search fails
                    if len(results) == 0:  # Only if we haven't found any results yet
                        text_filter = {
                            '$or': [
                                {'name': {'$regex': search_term, '$options': 'i'}},
                                {'text_content': {'$regex': search_term, '$options': 'i'}},
                                {'tags': {'$in': [search_term]}}
                            ]
                        }
                        
                        # Combine with existing filters
                        final_filter = {'$and': [mongo_filter, text_filter]} if mongo_filter else text_filter
                        
                        # Query MongoDB with text search
                        docs = await db.documents.find(final_filter, {"_id": 0}).limit(query.limit).to_list(query.limit)
                        
                        for doc in docs:
                            results.append(DocumentResponse(
                                id=doc['id'],
                                name=doc['name'],
                                file_type=doc['file_type'],
                                size=doc['size'],
                                tags=doc['tags'],
                                file_path=doc['file_path'],
                                snippet=get_snippet(doc.get('text_content', '')),
                                created_at=doc['created_at'],
                                updated_at=doc['updated_at']
                            ))
        else:
            # Return documents with filters applied (or all if no filters)
            docs = await db.documents.find(mongo_filter, {"_id": 0}).sort('created_at', -1).limit(query.limit).to_list(query.limit)
            
            for doc in docs:
                results.append(DocumentResponse(
                    id=doc['id'],
                    name=doc['name'],
                    file_type=doc['file_type'],
                    size=doc['size'],
                    tags=doc['tags'],
                    file_path=doc['file_path'],
                    snippet=get_snippet(doc.get('text_content', '')),
                    created_at=doc['created_at'],
                    updated_at=doc['updated_at']
                ))
        
        return {"results": [r.model_dump() for r in results], "total": len(results)}
    
    except Exception as e:
        logging.error(f"Error searching documents: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/documents/{doc_id}")
async def get_document(doc_id: str):
    """Get document details"""
    doc = await db.documents.find_one({"id": doc_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    return doc

@api_router.get("/documents/{doc_id}/download")
async def download_document(doc_id: str):
    """Download document file"""
    doc = await db.documents.find_one({"id": doc_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    file_path = Path(doc['file_path'])
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    
    return FileResponse(file_path, filename=doc['name'])

@api_router.put("/documents/{doc_id}/tags")
async def update_tags(doc_id: str, request: UpdateTagsRequest):
    """Update document tags"""
    result = await db.documents.update_one(
        {"id": doc_id},
        {"$set": {"tags": request.tags, "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="Document not found")
    
    return {"message": "Tags updated successfully", "tags": request.tags}

@api_router.delete("/documents/{doc_id}")
async def delete_document(doc_id: str):
    """Delete document"""
    doc = await db.documents.find_one({"id": doc_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Delete file
    file_path = Path(doc['file_path'])
    if file_path.exists():
        file_path.unlink()
    
    # Delete from Chroma
    try:
        collection.delete(ids=[doc_id])
    except:
        pass
    
    # Delete from MongoDB
    await db.documents.delete_one({"id": doc_id})
    
    return {"message": "Document deleted successfully"}

@api_router.get("/stats")
async def get_stats():
    """Get statistics"""
    total_docs = await db.documents.count_documents({})
    
    # Get file type distribution
    pipeline = [
        {"$group": {"_id": "$file_type", "count": {"$sum": 1}}},
        {"$sort": {"count": -1}}
    ]
    file_types = await db.documents.aggregate(pipeline).to_list(100)
    
    # Get all unique tags
    all_docs = await db.documents.find({}, {"tags": 1, "_id": 0}).to_list(1000)
    all_tags = set()
    for doc in all_docs:
        all_tags.update(doc.get('tags', []))
    
    return {
        "total_documents": total_docs,
        "file_types": file_types,
        "tags": sorted(list(all_tags))
    }

# Include the router in the main app
app.include_router(api_router)

# Configure CORS
origins = os.environ.get('CORS_ORIGINS', 'http://localhost:3000,http://127.0.0.1:3000').split(',')

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()

# Add health check endpoint
@app.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now(timezone.utc).isoformat()}
